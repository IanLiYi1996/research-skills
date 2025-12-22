#!/usr/bin/env python3
"""
PPTX 组装器

将生成的图像组装成 PowerPoint 演示文稿。
"""

import argparse
import sys
from pathlib import Path
from typing import List, Tuple
import re

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    print("错误: python-pptx 未安装")
    print("安装方法: pip install python-pptx")
    sys.exit(1)

try:
    from PIL import Image
except ImportError:
    print("错误: Pillow 未安装")
    print("安装方法: pip install Pillow")
    sys.exit(1)


def get_image_files(images_dir: str, sort_by: str = 'name') -> List[Path]:
    """
    获取图像目录中的所有图像文件

    Args:
        images_dir: 图像目录路径
        sort_by: 排序方式 ('name' 或 'date')

    Returns:
        排序后的图像文件列表
    """
    images_path = Path(images_dir)

    if not images_path.exists():
        raise FileNotFoundError(f"图像目录不存在: {images_dir}")

    # 支持的图像格式
    extensions = ['.png', '.jpg', '.jpeg', '.webp']

    # 获取所有图像文件
    image_files = []
    for ext in extensions:
        image_files.extend(images_path.glob(f'*{ext}'))
        image_files.extend(images_path.glob(f'*{ext.upper()}'))

    if not image_files:
        raise ValueError(f"在 {images_dir} 中未找到任何图像文件")

    # 排序
    if sort_by == 'name':
        # 尝试提取数字进行自然排序
        def extract_number(path: Path) -> int:
            match = re.search(r'(\d+)', path.stem)
            return int(match.group(1)) if match else 0

        image_files.sort(key=extract_number)
    elif sort_by == 'date':
        image_files.sort(key=lambda p: p.stat().st_mtime)
    else:
        raise ValueError(f"不支持的排序方式: {sort_by}")

    return image_files


def get_image_dimensions(image_path: Path) -> Tuple[int, int]:
    """
    获取图像尺寸

    Args:
        image_path: 图像文件路径

    Returns:
        (width, height) 元组
    """
    with Image.open(image_path) as img:
        return img.size


def create_pptx_from_images(
    image_files: List[Path],
    output_file: str,
    layout: str = 'full',
    add_textbox: bool = False,
    textbox_position: str = 'bottom'
) -> None:
    """
    从图像创建 PPTX

    Args:
        image_files: 图像文件列表
        output_file: 输出 PPTX 文件路径
        layout: 布局方式 ('full' 铺满, 'fit' 适应)
        add_textbox: 是否添加可编辑文本框
        textbox_position: 文本框位置 ('top', 'bottom', 'left', 'right')
    """
    print("="*60)
    print("PPTX 组装器")
    print("="*60 + "\n")

    # 创建演示文稿
    prs = Presentation()

    # 设置幻灯片尺寸为 16:9
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    print(f"步骤 1: 创建演示文稿")
    print("-" * 40)
    print(f"✓ 幻灯片尺寸: 16:9 (10\" × 5.625\")\n")

    # 添加幻灯片
    print(f"步骤 2: 添加幻灯片图像")
    print("-" * 40)

    for i, image_file in enumerate(image_files, 1):
        # 使用空白布局
        blank_layout = prs.slide_layouts[6]  # 6 是空白布局
        slide = prs.slides.add_slide(blank_layout)

        # 获取图像尺寸
        img_width, img_height = get_image_dimensions(image_file)
        img_aspect = img_width / img_height

        # 计算放置尺寸和位置
        slide_aspect = prs.slide_width / prs.slide_height

        if layout == 'full':
            # 铺满模式：图像填满整个幻灯片
            left = Inches(0)
            top = Inches(0)
            width = prs.slide_width
            height = prs.slide_height

        elif layout == 'fit':
            # 适应模式：保持图像宽高比，适应幻灯片
            if img_aspect > slide_aspect:
                # 图像更宽，以宽度为准
                width = prs.slide_width
                height = width / img_aspect
                left = Inches(0)
                top = (prs.slide_height - height) / 2
            else:
                # 图像更高，以高度为准
                height = prs.slide_height
                width = height * img_aspect
                top = Inches(0)
                left = (prs.slide_width - width) / 2
        else:
            raise ValueError(f"不支持的布局方式: {layout}")

        # 添加图像
        slide.shapes.add_picture(
            str(image_file),
            left, top,
            width=width,
            height=height
        )

        print(f"  [{i}/{len(image_files)}] {image_file.name}")

        # 可选：添加文本框
        if add_textbox:
            # 文本框尺寸
            tb_width = Inches(8)
            tb_height = Inches(0.8)

            # 根据位置计算坐标
            if textbox_position == 'bottom':
                tb_left = (prs.slide_width - tb_width) / 2
                tb_top = prs.slide_height - tb_height - Inches(0.3)
            elif textbox_position == 'top':
                tb_left = (prs.slide_width - tb_width) / 2
                tb_top = Inches(0.3)
            elif textbox_position == 'left':
                tb_left = Inches(0.3)
                tb_top = (prs.slide_height - tb_height) / 2
            elif textbox_position == 'right':
                tb_left = prs.slide_width - tb_width - Inches(0.3)
                tb_top = (prs.slide_height - tb_height) / 2
            else:
                tb_left = Inches(1)
                tb_top = prs.slide_height - Inches(1)

            # 添加文本框
            txBox = slide.shapes.add_textbox(tb_left, tb_top, tb_width, tb_height)
            tf = txBox.text_frame
            tf.text = ""  # 空文本框，供编辑

            # 设置文本格式
            p = tf.paragraphs[0]
            p.font.size = Pt(14)
            p.font.name = '微软雅黑'

    print(f"\n✓ 已添加 {len(image_files)} 张幻灯片\n")

    # 保存
    print(f"步骤 3: 保存演示文稿")
    print("-" * 40)

    output_path = Path(output_file)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    prs.save(str(output_path))

    print(f"✓ 演示文稿已保存到: {output_file}")
    print(f"  - 幻灯片数量: {len(prs.slides)}")
    print(f"  - 文件大小: {output_path.stat().st_size / 1024 / 1024:.2f} MB")

    print("\n" + "="*60)
    print("✓ PPTX 组装完成！")
    print("="*60)
    print(f"\n📊 可以用 PowerPoint 或 LibreOffice 打开:")
    print(f"  {output_file}")


def main():
    parser = argparse.ArgumentParser(
        description='从图像组装 PowerPoint 演示文稿',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
示例:
  # 基本用法
  %(prog)s --images ./slide-images/ --output presentation.pptx

  # 添加可编辑文本框
  %(prog)s --images ./slide-images/ --output presentation.pptx --add-textbox

  # 指定文本框位置
  %(prog)s --images ./images/ --output ppt.pptx --add-textbox --textbox-position top

  # 使用适应布局（保持图像比例）
  %(prog)s --images ./images/ --output ppt.pptx --layout fit

  # 按日期排序
  %(prog)s --images ./images/ --output ppt.pptx --sort date

布局说明:
  full  - 图像铺满整个幻灯片（可能裁剪，默认）
  fit   - 图像适应幻灯片（保持比例，可能有边距）

文本框位置:
  top, bottom, left, right
        """
    )

    parser.add_argument(
        '--images', '-i',
        required=True,
        help='图像目录路径'
    )

    parser.add_argument(
        '--output', '-o',
        default='presentation.pptx',
        help='输出 PPTX 文件路径（默认: presentation.pptx）'
    )

    parser.add_argument(
        '--layout', '-l',
        default='full',
        choices=['full', 'fit'],
        help='布局方式（默认: full）'
    )

    parser.add_argument(
        '--add-textbox',
        action='store_true',
        help='为每页添加可编辑文本框'
    )

    parser.add_argument(
        '--textbox-position',
        default='bottom',
        choices=['top', 'bottom', 'left', 'right'],
        help='文本框位置（默认: bottom）'
    )

    parser.add_argument(
        '--sort',
        default='name',
        choices=['name', 'date'],
        help='图像排序方式（默认: name）'
    )

    args = parser.parse_args()

    try:
        # 获取图像文件
        image_files = get_image_files(args.images, args.sort)

        if not image_files:
            print(f"❌ 错误: 在 {args.images} 中未找到图像文件")
            sys.exit(1)

        print(f"✓ 找到 {len(image_files)} 个图像文件\n")

        # 创建 PPTX
        create_pptx_from_images(
            image_files=image_files,
            output_file=args.output,
            layout=args.layout,
            add_textbox=args.add_textbox,
            textbox_position=args.textbox_position
        )

    except KeyboardInterrupt:
        print("\n\n⚠️  用户中断")
        sys.exit(1)
    except Exception as e:
        print(f"\n❌ 错误: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)


if __name__ == "__main__":
    main()
